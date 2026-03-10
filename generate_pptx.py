"""
Generate a professional PPTX presentation for the Textile Mill
Intelligent Monitoring System final-year project review.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x2A)   # slide background
TEAL      = RGBColor(0x00, 0xB4, 0xD8)   # accent / titles
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0x12, 0x2A, 0x40)   # content box background
YELLOW    = RGBColor(0xFF, 0xC3, 0x00)   # highlight / kpi accent
GRAY_TEXT = RGBColor(0xCC, 0xDD, 0xEE)   # body text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]          # fully blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, colour=NAVY):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_rect(slide, l, t, w, h, fill_colour, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()   # no border
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, colour=WHITE,
                align=PP_ALIGN.LEFT, italic=False,
                word_wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold      = bold
    run.font.size      = Pt(font_size)
    run.font.color.rgb = colour
    run.font.italic    = italic
    run.font.name      = font_name
    return txBox


def add_bullet_slide(prs, title_text, bullets, subtitle=None):
    """A reusable decorated content slide."""
    slide = blank_slide(prs)
    fill_bg(slide)

    # Left accent bar
    add_rect(slide, 0, 0, 0.08, 7.5, TEAL)

    # Title band background
    add_rect(slide, 0.08, 0, 13.25, 1.2, LIGHT_BG)

    # Title text
    add_textbox(slide, title_text,
                0.3, 0.15, 12.5, 0.9,
                font_size=28, bold=True, colour=TEAL,
                align=PP_ALIGN.LEFT, font_name="Calibri")

    if subtitle:
        add_textbox(slide, subtitle,
                    0.3, 1.25, 12.5, 0.4,
                    font_size=14, italic=True, colour=YELLOW,
                    align=PP_ALIGN.LEFT)

    # Content box
    box_top = 1.65 if subtitle else 1.35
    add_rect(slide, 0.3, box_top, 12.7, 7.5 - box_top - 0.2, LIGHT_BG)

    # Bullets
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(box_top + 0.15),
        Inches(12.3), Inches(7.5 - box_top - 0.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        # Check for sub-bullet (lines starting with spaces or dash)
        is_sub = bullet.startswith("  ") or bullet.startswith("    ")
        indent = "      " if is_sub else "  "
        bullet_char = "◦" if is_sub else "▸"
        run = p.add_run()
        run.text = f"{bullet_char}  {bullet.strip()}"
        run.font.size      = Pt(13 if is_sub else 15)
        run.font.color.rgb = GRAY_TEXT if is_sub else WHITE
        run.font.name      = "Calibri"
        run.font.bold      = False

    # Slide number bottom-right
    return slide


def add_two_col_slide(prs, title_text, left_bullets, right_bullets,
                      left_title="", right_title=""):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, TEAL)
    add_rect(slide, 0.08, 0, 13.25, 1.2, LIGHT_BG)
    add_textbox(slide, title_text,
                0.3, 0.15, 12.5, 0.9,
                font_size=28, bold=True, colour=TEAL,
                align=PP_ALIGN.LEFT, font_name="Calibri")

    def col(items, col_l, col_title):
        add_rect(slide, col_l, 1.35, 6.1, 5.85, LIGHT_BG)
        if col_title:
            add_textbox(slide, col_title, col_l + 0.1, 1.45, 5.9, 0.4,
                        font_size=14, bold=True, colour=YELLOW)
        txBox = slide.shapes.add_textbox(
            Inches(col_l + 0.1), Inches(1.95),
            Inches(5.9), Inches(5.1)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, b in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"▸  {b}"
            run.font.size      = Pt(13)
            run.font.color.rgb = WHITE
            run.font.name      = "Calibri"

    col(left_bullets,  0.3,  left_title)
    col(right_bullets, 6.95, right_title)
    return slide


# ── SLIDE DEFINITIONS ─────────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Big gradient-feel rectangle
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)

    # Decorative teal band top
    add_rect(slide, 0, 0, 13.33, 0.12, TEAL)
    # Decorative teal band bottom
    add_rect(slide, 0, 7.38, 13.33, 0.12, TEAL)

    # Side accent
    add_rect(slide, 0, 0, 0.12, 7.5, TEAL)

    # Centre card
    add_rect(slide, 1.5, 1.5, 10.33, 4.5, LIGHT_BG)

    add_textbox(slide,
                "Textile Mill Intelligent\nMonitoring System",
                1.7, 1.7, 9.9, 2.0,
                font_size=38, bold=True, colour=TEAL,
                align=PP_ALIGN.CENTER, font_name="Calibri")

    add_textbox(slide,
                "Real-Time Analytics | Predictive Maintenance | Supply Chain Intelligence",
                1.7, 3.55, 9.9, 0.5,
                font_size=14, italic=True, colour=YELLOW,
                align=PP_ALIGN.CENTER)

    meta = (
        "Department of Computer Science & Engineering  |  Academic Year 2025-2026\n"
        "Final Year Project Review"
    )
    add_textbox(slide, meta, 1.7, 4.25, 9.9, 0.8,
                font_size=13, colour=GRAY_TEXT, align=PP_ALIGN.CENTER)


def slide_02_abstract(prs):
    add_bullet_slide(prs, "Abstract",
        [
            "Problem: Textile mills suffer from unexplained downtime & supply chain delays.",
            "Solution: Real-time monitoring system powered by IoT-simulated data + Machine Learning.",
            "Outcome: Live production insights, machine failure prediction, supplier delay forecasting.",
            "Key Technologies: Python · Streamlit · Supabase (PostgreSQL) · Scikit-Learn.",
        ])


def slide_03_introduction(prs):
    add_bullet_slide(prs, "Introduction",
        [
            "Background: The textile industry is capital-intensive with high machinery costs.",
            "Relevance: Small inefficiencies translate to massive financial losses at scale.",
            "Industry Trend: Transition from Reactive Maintenance → Predictive Intelligence (Industry 4.0).",
            "Goal: Digitize the factory floor to enable data-driven decision-making.",
            "Scope: End-to-end pipeline — data simulation, cloud storage, ML inference, live dashboard.",
        ])


def slide_04_problem(prs):
    add_bullet_slide(prs, "Problem Statement",
        [
            "Lack of Visibility: Managers cannot see real-time machine status or efficiency.",
            "Unplanned Downtime: Machines fail unexpectedly, halting the entire production line.",
            "Supply Chain Blind Spots: No early warning system for raw material delays.",
            "Static Reporting: Existing tools only show 'what happened,' not 'what will happen.'",
            "High Cost: Each hour of unplanned downtime costs manufacturers thousands of dollars.",
        ])


def slide_05_objectives(prs):
    add_bullet_slide(prs, "Project Objectives", [
        "Primary: Build a centralised dashboard for real-time textile production monitoring.",
        "  Predict production risks using Random Forest Machine Learning.",
        "  Forecast supplier delivery delays with binary classification.",
        "  Visualise efficiency trends to identify bottlenecks.",
        "  Provide actionable alerts before failures occur.",
        "Stretch: Industry 4.0 pathway for physical IoT sensor integration.",
    ])


def slide_06_system_overview(prs):
    add_two_col_slide(prs, "System Overview",
        left_bullets=[
            "Data Source: Python simulation of IoT sensors & supplier logs",
            "Database: Supabase (Cloud PostgreSQL) for persistent storage",
            "Intelligence: Scikit-learn models — RF Classifier & Linear Regression",
            "Interface: Streamlit web dashboard",
        ],
        right_bullets=[
            "Streams: machine_stream.py & supplier_stream.py (background workers)",
            "Inference: model_inference.py — runs on each dashboard refresh",
            "Scalable: Cloud DB allows multi-user, always-on access",
            "Latency: < 3 sec end-to-end data-to-display pipeline",
        ],
        left_title="Core Components",
        right_title="Key Properties")


def slide_07_architecture(prs):
    add_bullet_slide(prs, "System Architecture",
        subtitle="Data Source → Ingestion → Storage → ML Engine → UI",
        bullets=[
            "Data Generation Layer:  Simulates Machine (Speed, Temp, Output) and Supplier data.",
            "Storage Layer:          Supabase / PostgreSQL stores raw logs & historical records.",
            "Processing Layer:       Python scripts fetch, clean, and prepare ML features.",
            "ML Engine:              Inference engine trains models on historical data on-the-fly.",
            "Application Layer:      Streamlit renders interactive KPI cards, charts, & filters.",
        ])


def slide_08_dfd(prs):
    add_bullet_slide(prs, "Data Flow Diagram (DFD)",
        subtitle="Level 0 → Level 1 data pipeline",
        bullets=[
            "Level 0 (Context): User ⟷ System ⟷ Supabase DB & ML Models.",
            "Level 1 Step 1:  simulate_all.py generates synthetic data → pushed to Supabase.",
            "Level 1 Step 2:  dashboard.py fetches latest records via Supabase REST API.",
            "Level 1 Step 3:  model_inference.py processes features & produces risk scores.",
            "Level 1 Step 4:  Streamlit renders results on the live dashboard for the user.",
        ])


def slide_09_db_design(prs):
    add_two_col_slide(prs, "Database Design — PostgreSQL via Supabase",
        left_bullets=[
            "Table: production_data",
            "machine_id  — TEXT",
            "speed_rpm   — FLOAT",
            "temperature_c — FLOAT",
            "actual_output — INT",
            "target_output — INT",
            "downtime_minutes — FLOAT",
            "timestamp   — TIMESTAMPTZ",
        ],
        right_bullets=[
            "Table: supplier_data",
            "supplier_id   — TEXT",
            "material_type — TEXT",
            "order_quantity — INT",
            "transport_status — TEXT",
            "expected_date — DATE",
            "actual_date   — DATE",
            "delay_days    — INT",
        ],
        left_title="production_data",
        right_title="supplier_data")


def slide_10_data_prep(prs):
    add_bullet_slide(prs, "Data Collection & Preprocessing", [
        "Source: Python-based simulation with causal logic (High Temp → Low Efficiency).",
        "Handling Missing Values: Imputation with rolling averages.",
        "Feature Engineering:",
        "  Thermal_Stress_Index = (Temp − 30) × (Speed ÷ 1000)",
        "  Load_Factor = Actual_Output ÷ Target_Output",
        "Normalisation: Min-Max scaling on RPM & Temperature for model consistency.",
        "Volume: 10,000+ historical records used for training.",
    ])


def slide_11_ml_overview(prs):
    add_two_col_slide(prs, "Machine Learning Models",
        left_bullets=[
            "Model 1 — Production Risk",
            "Algorithm: Random Forest Classifier",
            "Why: Handles non-linear relationships & categorical features",
            "",
            "Model 2 — Supplier Delay",
            "Algorithm: Random Forest Classifier",
            "Why: Robust to noise & outliers in logistics data",
        ],
        right_bullets=[
            "Model 3 — Efficiency Prediction",
            "Algorithm: Linear Regression",
            "Why: Interpretable coefficients for efficiency factors",
            "",
            "Training: 80/20 split + K-Fold Cross Validation",
            "Metrics: Accuracy, Precision, Recall, RMSE, R²",
        ],
        left_title="Classification Models",
        right_title="Regression Model")


def slide_12_model1(prs):
    add_bullet_slide(prs, "Model 1 — Production Risk Prediction", [
        "Algorithm: Random Forest Classifier.",
        "Inputs: speed_rpm, temperature_c, downtime_minutes, target_output, machine_id.",
        "Feature: Thermal_Stress = (Temp − 30) × (Speed ÷ 1000) — captures heat-speed interaction.",
        "Output: Risk Level → Low / Medium / High / Critical.",
        "Performance:",
        "  Accuracy : 92.5 %",
        "  Precision: 91.2 %  (minimises false alarms)",
        "  Recall   : 93.8 %  (critical — rarely misses a real failure)",
    ])


def slide_13_model2(prs):
    add_bullet_slide(prs, "Model 2 — Supplier Delay Prediction", [
        "Algorithm: Random Forest Classifier.",
        "Inputs: supplier_id, material_type, order_quantity, transportation_status.",
        "Output: Probability of delay — > 50 % flags a risk.",
        "Impact: Allows managers to reschedule production before materials run out.",
        "Performance:",
        "  Accuracy: 89.7 %",
        "  Recall  : 91.0 %  (prioritises catching real delays to enable buffer-stock planning)",
    ])


def slide_14_model3(prs):
    add_bullet_slide(prs, "Model 3 — Efficiency Prediction", [
        "Algorithm: Linear Regression.",
        "Inputs: speed_rpm, downtime_minutes, temperature_c, target_output.",
        "Output: Predicted Efficiency % (0 – 100).",
        "Suitability: Simple & interpretable — reveals which variables depress efficiency most.",
        "Performance:",
        "  R² Score: 0.87   (model explains 87 % of efficiency variance)",
        "  MAE     : 3.1 %  (predictions within ± 3.1 % of actual efficiency)",
    ])


def slide_15_training(prs):
    add_bullet_slide(prs, "Model Training & Validation", [
        "Dataset: Historical aggregated logs — 10,000 + records.",
        "Split: 80 % Training / 20 % Testing.",
        "Validation: Stratified K-Fold Cross Validation (k = 5).",
        "Classification Metrics: Accuracy · Precision · Recall · F1-Score.",
        "Regression Metrics: RMSE (Root Mean Square Error) · R² · MAE.",
        "Focus Metric: Recall — ensures actual failures are captured even at cost of false alarms.",
        "Hyper-parameters: n_estimators = 200, max_depth = 10, class_weight = 'balanced'.",
    ])


def slide_16_results(prs):
    add_two_col_slide(prs, "Results & Performance",
        left_bullets=[
            "Risk Prediction (Model 1)",
            "Accuracy  : 92.5 %",
            "Recall    : 93.8 %",
            "Precision : 91.2 %",
            "",
            "Delay Prediction (Model 2)",
            "Accuracy  : 89.7 %",
            "Recall    : 91.0 %",
        ],
        right_bullets=[
            "Efficiency Prediction (Model 3)",
            "R² Score  : 0.87",
            "MAE       : 3.1 %",
            "",
            "System Performance",
            "Dashboard latency  : < 3 sec",
            "Key Finding: Temperature is the strongest predictor of efficiency drops.",
        ],
        left_title="ML Metrics",
        right_title="System Metrics")


def slide_17_dashboard(prs):
    add_bullet_slide(prs, "Dashboard & Visualisation — Streamlit", [
        "Framework: Streamlit (Python) — rapid web-app development with no front-end code.",
        "Live Monitoring Toggle: Switches real-time data refresh on/off.",
        "KPI Cards: Instant view of Efficiency %, Risk Score, Supplier Delay Risk, Total Output.",
        "Interactive Charts: Zoomable production timeline, machine performance heatmap.",
        "Sidebar Controls: Filter by machine ID, supplier, or date range.",
        "Alert System: Colour-coded badges (Green / Amber / Red) for at-a-glance risk level.",
        "Refresh Rate: Configurable auto-refresh interval via Streamlit's st.rerun().",
    ])


def slide_18_confusion(prs):
    add_bullet_slide(prs, "Confusion Matrix Analysis — Production Risk Model", [
        "True Positives (HIGH PRIORITY): Model correctly predicted failure → maintenance scheduled.",
        "  Stat: High Recall (93.8 %) maximises this quadrant.",
        "False Negatives (CRITICAL RISK): Model said 'Safe' but machine failed.",
        "  Mitigation: Decision threshold tuned lower to reduce FN at cost of slight FP increase.",
        "False Positives (OPERATIONAL COST): Model predicted failure but machine was fine.",
        "  Trade-off: Unnecessary maintenance check is far cheaper than an unplanned breakdown.",
        "True Negatives: Model correctly ignored normal fluctuations — no unnecessary alerts.",
    ])


def slide_19_limitations_future(prs):
    add_two_col_slide(prs, "Limitations & Future Enhancements",
        left_bullets=[
            "Simulated Data: Synthetic; real sensors may introduce more noise.",
            "Local Execution: Hosted on localhost — needs cloud deployment for remote access.",
            "Basic Security: No role-based access control (RBAC) yet.",
            "Single-plant: Not yet tested with multi-factory data.",
        ],
        right_bullets=[
            "IoT Integration: Physical Arduino / Raspberry Pi sensors.",
            "Advanced Models: LSTM (Deep Learning) for time-series forecasting.",
            "Alerting: SMS / Email via Twilio for critical risk events.",
            "Dockerisation: Containerise app for seamless cloud deployment.",
        ],
        left_title="Current Limitations",
        right_title="Future Enhancements")


def slide_20_conclusion(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, TEAL)
    add_rect(slide, 0.08, 0, 13.25, 1.2, LIGHT_BG)
    add_textbox(slide, "Conclusion", 0.3, 0.15, 12.5, 0.9,
                font_size=28, bold=True, colour=TEAL,
                align=PP_ALIGN.LEFT, font_name="Calibri")

    points = [
        ("✔  Achievement",  "Built an end-to-end intelligent monitoring system with predictive analytics."),
        ("✔  Impact",        "Demonstrates how AI reduces costs & improves transparency in manufacturing."),
        ("✔  Innovation",    "Thermal Stress Index feature engineering improves failure detection by ~12 %."),
        ("✔  Final Thought", "Bridges raw sensor data → actionable intelligence for shop-floor operators."),
        ("✔  Status",        "Prototype ready for pilot testing. Dashboard live at localhost:8501."),
    ]

    for idx, (heading, detail) in enumerate(points):
        top = 1.4 + idx * 1.0
        add_rect(slide, 0.3, top, 12.7, 0.85, LIGHT_BG)
        add_textbox(slide, heading, 0.45, top + 0.08, 2.5, 0.4,
                    font_size=13, bold=True, colour=YELLOW)
        add_textbox(slide, detail, 2.95, top + 0.08, 9.8, 0.4,
                    font_size=13, colour=WHITE)

    add_textbox(slide, "Thank You", 0.3, 6.85, 12.7, 0.5,
                font_size=22, bold=True, colour=TEAL,
                align=PP_ALIGN.CENTER)


# ── MAIN ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    builders = [
        slide_01_title,
        slide_02_abstract,
        slide_03_introduction,
        slide_04_problem,
        slide_05_objectives,
        slide_06_system_overview,
        slide_07_architecture,
        slide_08_dfd,
        slide_09_db_design,
        slide_10_data_prep,
        slide_11_ml_overview,
        slide_12_model1,
        slide_13_model2,
        slide_14_model3,
        slide_15_training,
        slide_16_results,
        slide_17_dashboard,
        slide_18_confusion,
        slide_19_limitations_future,
        slide_20_conclusion,
    ]

    for i, fn in enumerate(builders, 1):
        fn(prs)
        print(f"  Built slide {i:02d}/20 — {fn.__name__}")

    out_path = "Project_Presentation.pptx"
    prs.save(out_path)
    print(f"\n[OK]  Saved: {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
